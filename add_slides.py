# Copyright (c) 2026 Sebastian Herrera Betancur
# Biomicrosystems Research Group | Universidad de los Andes
# PROPRIETARY CODE - Unauthorized use, copying or distribution is strictly prohibited.

from pptx import Presentation
from pptx.util import Inches, Pt
def add_slides(pptx_path):
    try:
        prs = Presentation(pptx_path)
    except Exception:
        prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Arquitectura de Hardware"
    tf = content.text_frame
    tf.text = "Integración de Sensores y Control"
    p = tf.add_paragraph()
    p.text = "• AS726X: Espectrofotómetro de 6 canales (Visible)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• MLX90640: Cámara térmica (Arreglo focal 24x32)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Comunicación: Dual I2C (Bus 1 y Bus 22) + GPIO 17 (LED)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Modo Emulación: Soporte completo para desarrollo en Windows."
    p.level = 1
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Arquitectura de Software (Backend)"
    tf = content.text_frame
    tf.text = "Potencia y Rigor Analítico"
    p = tf.add_paragraph()
    p.text = "• Django 6.0: Servidor web y API robusta."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Motor Matemático: NumPy y SciPy para procesamiento de señales."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Algoritmos: Regresión OLS, Shapiro-Wilk y pruebas t-Student."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Estado Global: Gestión thread-safe de sesión instrumental única."
    p.level = 1
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Interfaz de Usuario (Frontend)"
    tf = content.text_frame
    tf.text = "Experiencia de Laboratorio Moderna"
    p = tf.add_paragraph()
    p.text = "• SPA (Single Page Application): Vanilla JS de alta velocidad."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Visualización: Chart.js interactivo (Espectros y Termografía)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• i18n Dinámica: Cambio instantáneo de idioma (EN/ES)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• UX: Modo Oscuro/Claro basado en variables CSS nativas."
    p.level = 1
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Gestión de Datos e Integridad"
    tf = content.text_frame
    tf.text = "Trazabilidad Científica"
    p = tf.add_paragraph()
    p.text = "• Modelos: Beer-Lambert, Calibración y Transferencia."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Exportación: CSV enriquecido con metadatos y frames térmicos."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Métricas: Cálculo automático de LOD, LOQ y RMSE."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Seguridad: Sincronización térmica por punto de medición."
    p.level = 1
    prs.save(pptx_path)
    print(f"Slides added to {pptx_path}")
if __name__ == "__main__":
    add_slides('Presentacion_Espectrofotometro.pptx')